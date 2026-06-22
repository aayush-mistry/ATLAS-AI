import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()

# Set 16:9 widescreen layout
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette (Premium Dark Theme matching Cyber-NOC)
BG_COLOR = RGBColor(11, 11, 14)       # Ultra dark blue-gray
CARD_COLOR = RGBColor(20, 20, 26)     # Slightly lighter dark card
ACCENT_BLUE = RGBColor(59, 130, 246)  # Electric Blue
ACCENT_PURPLE = RGBColor(147, 51, 234)# Neon Purple
ACCENT_GREEN = RGBColor(16, 185, 129) # Success Green
TEXT_WHITE = RGBColor(255, 255, 255)  # Pure White
TEXT_MUTED = RGBColor(156, 163, 175)  # Muted Gray
TEXT_BODY = RGBColor(209, 213, 219)   # Light Gray

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_header(slide, title_text, category_text="ATLAS // AUTONOMOUS LEGAL ENTITY"):
    # Category Tracker (Top-Left small text)
    txBox_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
    tf_cat = txBox_cat.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = "Segoe UI"
    p_cat.font.size = Pt(9)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_BLUE

    # Main Title
    txBox_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf_title = txBox_title.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE

def add_footer(slide, current_slide, total_slides=15):
    # Bottom page number & footer
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = f"ATLAS AI  |  HACKATHON SUBMISSION  |  SLIDE {current_slide} OF {total_slides}"
    p.font.name = "Segoe UI"
    p.font.size = Pt(8)
    p.font.color.rgb = TEXT_MUTED

def add_card(slide, left, top, width, height, title="", border_color=None):
    # Draw a card background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_COLOR
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background() # borderless
        
    if title:
        txBox = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

# ----------------- SLIDE 1: COVER SLIDE -----------------
slide_layout = prs.slide_layouts[6] # Blank
slide1 = prs.slides.add_slide(slide_layout)
apply_background(slide1)

# Large Logo / Title
txBox = slide1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(3.0))
tf = txBox.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "ATLAS // AI"
p1.font.name = "Segoe UI"
p1.font.size = Pt(64)
p1.font.bold = True
p1.font.color.rgb = TEXT_WHITE
p1.space_after = Pt(10)

p2 = tf.add_paragraph()
p2.text = "Autonomous Legal Entities Operating on Decentralized Rails"
p2.font.name = "Segoe UI"
p2.font.size = Pt(22)
p2.font.color.rgb = ACCENT_BLUE
p2.space_after = Pt(20)

p3 = tf.add_paragraph()
p3.text = "Next-Gen AI Business Operations Command Center with Smart Contract Governance"
p3.font.name = "Segoe UI"
p3.font.size = Pt(14)
p3.font.color.rgb = TEXT_MUTED

# Author info / bottom info
txBox_author = slide1.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(5.0), Inches(1.0))
tf_author = txBox_author.text_frame
p_auth = tf_author.paragraphs[0]
p_auth.text = "DEVELOPED BY: TEAM MERIDIAN"
p_auth.font.name = "Segoe UI"
p_auth.font.size = Pt(10)
p_auth.font.bold = True
p_auth.font.color.rgb = TEXT_WHITE

p_auth2 = tf_author.add_paragraph()
p_auth2.text = "Hack Aarambh 2026 Submission"
p_auth2.font.name = "Segoe UI"
p_auth2.font.size = Pt(10)
p_auth2.font.color.rgb = TEXT_MUTED
add_footer(slide1, 1)

# ----------------- SLIDE 2: THE PROBLEM -----------------
slide2 = prs.slides.add_slide(slide_layout)
apply_background(slide2)
add_header(slide2, "The Traditional Business Friction")

# Left Box: Problem statements
txBox_left = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5))
tf_left = txBox_left.text_frame
tf_left.word_wrap = True

p = tf_left.paragraphs[0]
p.text = "The Modern Enterprise Bottleneck"
p.font.name = "Segoe UI"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE
p.space_after = Pt(15)

bullet_points = [
    ("Administrative Overhead", "Corporate structure setups require expensive legal consulting, notary validation, and endless paperwork filing."),
    ("Static Decision Loops", "Traditional managers react to market shifts slowly. Weather anomalies, supply disruptions, or demand drops cause delayed responses."),
    ("Siloed Financials & Payouts", "Settlements with workers take days/weeks. Bookkeeping is manually tracked in fragmented spreadsheets prone to errors."),
    ("Lack of Transparency", "Internal governance decisions are locked behind closed doors, providing no audit trail for stakeholders.")
]

for title, desc in bullet_points:
    p_title = tf_left.add_paragraph()
    p_title.text = f"• {title}"
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(14)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    p_title.space_before = Pt(8)
    
    p_desc = tf_left.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = "Segoe UI"
    p_desc.font.size = Pt(11)
    p_desc.font.color.rgb = TEXT_BODY
    p_desc.space_after = Pt(8)

# Right Box: Large callout
add_card(slide2, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5), "The Operational Cost of Friction", border_color=ACCENT_PURPLE)
txBox_right = slide2.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5.1), Inches(3.5))
tf_right = txBox_right.text_frame
tf_right.word_wrap = True

p_callout = tf_right.paragraphs[0]
p_callout.text = "Organizations spend up to 35% of revenue on administrative compliance, management overhead, and middle-man contract enforcement."
p_callout.font.name = "Segoe UI"
p_callout.font.size = Pt(16)
p_callout.font.color.rgb = TEXT_BODY
p_callout.space_after = Pt(25)

p_stat = tf_right.add_paragraph()
p_stat.text = "10x"
p_stat.font.name = "Segoe UI"
p_stat.font.size = Pt(54)
p_stat.font.bold = True
p_stat.font.color.rgb = ACCENT_PURPLE
p_stat.alignment = PP_ALIGN.CENTER

p_stat_lbl = tf_right.add_paragraph()
p_stat_lbl.text = "Slowdown in executing legal & vendor agreements"
p_stat_lbl.font.name = "Segoe UI"
p_stat_lbl.font.size = Pt(11)
p_stat_lbl.font.color.rgb = TEXT_MUTED
p_stat_lbl.alignment = PP_ALIGN.CENTER

add_footer(slide2, 2)

# ----------------- SLIDE 3: THE SOLUTION -----------------
slide3 = prs.slides.add_slide(slide_layout)
apply_background(slide3)
add_header(slide3, "ATLAS AI: The First Fully Autonomous Legal Entity")

# 3 Reusable Cards
card_data = [
    ("Autonomous Execution", "Governed entirely by an AI agent (CEO) that evaluates market variables and automatically takes ledger actions without a human middle-man.", Inches(0.8), ACCENT_BLUE),
    ("Decentralized Payouts", "Uses Smart Contract mechanics to secure worker payouts instantly. Funds are released programmatically upon verification of job completion.", Inches(4.8), ACCENT_PURPLE),
    ("Dynamic Market Adaptation", "Connects to active external Oracles (simulated weather/price feeds) to recalibrate retail prices, supply orders, and job rates in real-time.", Inches(8.8), ACCENT_GREEN)
]

for title, desc, left, color in card_data:
    add_card(slide3, left, Inches(1.8), Inches(3.7), Inches(4.5), title, border_color=color)
    txBox = slide3.shapes.add_textbox(left + Inches(0.25), Inches(2.6), Inches(3.2), Inches(3.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.name = "Segoe UI"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_BODY
    p.space_after = Pt(20)

    p_action = tf.add_paragraph()
    p_action.text = "ACTIVE PROTOCOL"
    p_action.font.name = "Segoe UI"
    p_action.font.size = Pt(9)
    p_action.font.bold = True
    p_action.font.color.rgb = color

add_footer(slide3, 3)

# ----------------- SLIDE 4: THE CORE CONCEPT -----------------
slide4 = prs.slides.add_slide(slide_layout)
apply_background(slide4)
add_header(slide4, "How it Works: The Operational Loop")

# 4 Column Process Flow
process_steps = [
    ("1. Oracle Ingestion", "System captures live external telemetry inputs (e.g. weather, price spikes, traffic indices) and feeds it to the AI decision node."),
    ("2. AI CEO Evaluation", "AI analyzes input telemetry, references current treasury balances, checks remaining inventory, and formulates strategic decisions."),
    ("3. Smart Ledger Write", "Decisions are committed to the business state. Contracts automatically update prices or dispatch open workforce positions."),
    ("4. Automated Payouts", "Workers execute jobs, report proof of work, and receive instantaneous payment from the on-chain multi-sig wallet.")
]

for i, (title, desc) in enumerate(process_steps):
    left = Inches(0.8 + i * 2.95)
    add_card(slide4, left, Inches(1.8), Inches(2.7), Inches(4.5), title, border_color=ACCENT_BLUE if i == 0 else None)
    txBox = slide4.shapes.add_textbox(left + Inches(0.2), Inches(2.5), Inches(2.3), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.name = "Segoe UI"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_BODY
    
    # Visual arrow indicating flow
    if i < 3:
        arr_box = slide4.shapes.add_textbox(left + Inches(2.65), Inches(3.6), Inches(0.3), Inches(0.5))
        arr_tf = arr_box.text_frame
        p_arr = arr_tf.paragraphs[0]
        p_arr.text = "➔"
        p_arr.font.name = "Segoe UI"
        p_arr.font.size = Pt(20)
        p_arr.font.color.rgb = ACCENT_PURPLE

add_footer(slide4, 4)

# ----------------- SLIDE 5: THE DEMO PLAYGROUND -----------------
slide5 = prs.slides.add_slide(slide_layout)
apply_background(slide5)
add_header(slide5, "The Simulated Sandbox: Lemonade Stand")

add_card(slide5, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5), "A Fully Integrated Economic Environment")

txBox = slide5.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(11.1), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "To demonstrate ATLAS AI in a self-sustaining environment, we simulated a micro-economy: a Lemonade Stand managed entirely by AI."
p.font.name = "Segoe UI"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE
p.space_after = Pt(15)

sandbox_details = [
    ("AI Purchasing Decisions", "When ingredients (Lemons, Sugar, Cups, Ice) run low, the AI autonomously orders more from virtual wholesalers, calculating the unit cost against current treasury margins."),
    ("Recruitment & Job Board", "The AI creates job openings ('Make Lemonade', 'Manage Sales'). Human users submit worker credentials and digital wallet addresses to claim tasks."),
    ("Autonomous Sales & Dynamic Pricing", "The AI adjusts price per cup (e.g. ₹20 on hot sunny days vs. ₹10 on rainy days). Sales are computed based on dynamic customer demand indexes.")
]

for title, desc in sandbox_details:
    p_title = tf.add_paragraph()
    p_title.text = f"• {title}: "
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(13)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    p_title.space_before = Pt(6)

    p_desc = tf.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = "Segoe UI"
    p_desc.font.size = Pt(11)
    p_desc.font.color.rgb = TEXT_BODY
    p_desc.space_after = Pt(6)

add_footer(slide5, 5)

# ----------------- SLIDE 6: THE ARCHITECTURE -----------------
slide6 = prs.slides.add_slide(slide_layout)
apply_background(slide6)
add_header(slide6, "Decentralized & Intelligent Stack")

# Left Column: Tech Stack list
txBox_left = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5))
tf_left = txBox_left.text_frame
tf_left.word_wrap = True

p = tf_left.paragraphs[0]
p.text = "System Integration Layers"
p.font.name = "Segoe UI"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE
p.space_after = Pt(15)

stack = [
    ("AI Reasoner Interface", "Claude 3.5 Sonnet / local LLMs serving as the business brain, evaluating trade-offs, price changes, and worker recruitment strategy."),
    ("Local DB & ORM", "SQLite Database managed through Prisma ORM for tracking inventory levels, historical decisions, transactions, and logs."),
    ("Web3 Smart Contracts", "Hardhat development framework with ethers.js to simulate transaction settlements, payroll deposits, and verified state changes."),
    ("Telemetry Frontend", "Next.js 16 + React 19 + Tailwind CSS + Recharts visualizer. Powered by Turbopack for ultra-fast dashboard reloads.")
]

for title, desc in stack:
    p_title = tf_left.add_paragraph()
    p_title.text = f"✔ {title}"
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(13)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    p_title.space_before = Pt(6)
    
    p_desc = tf_left.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = "Segoe UI"
    p_desc.font.size = Pt(10)
    p_desc.font.color.rgb = TEXT_BODY
    p_desc.space_after = Pt(6)

# Right Column: Visual Diagram Cards
add_card(slide6, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5), "Architecture Dataflow", border_color=ACCENT_PURPLE)
txBox_right = slide6.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5.1), Inches(3.5))
tf_right = txBox_right.text_frame
tf_right.word_wrap = True

diagram_steps = [
    ("[1] Market Oracle Update", "Updates weather & pricing constants."),
    ("[2] LLM CEO Reasoning Engine", "Evaluates inventory & runs logic."),
    ("[3] Hardhat Mock Blockchain Ledger", "Dispatches on-chain transactions."),
    ("[4] Next.js Telemetry NOC Dashboard", "Displays telemetry in real-time.")
]

for title, desc in diagram_steps:
    p_d_title = tf_right.add_paragraph()
    p_d_title.text = title
    p_d_title.font.name = "Segoe UI"
    p_d_title.font.size = Pt(12)
    p_d_title.font.bold = True
    p_d_title.font.color.rgb = ACCENT_PURPLE
    p_d_title.space_before = Pt(8)
    
    p_d_desc = tf_right.add_paragraph()
    p_d_desc.text = desc
    p_d_desc.font.name = "Segoe UI"
    p_d_desc.font.size = Pt(10)
    p_d_desc.font.color.rgb = TEXT_BODY
    p_d_desc.space_after = Pt(4)

add_footer(slide6, 6)

# ----------------- SLIDE 7: CYBER-NOC OPERATIONS DASHBOARD -----------------
slide7 = prs.slides.add_slide(slide_layout)
apply_background(slide7)
add_header(slide7, "Cyber-NOC Dashboard: Operation Command Center")

add_card(slide7, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5), "Premium Telemetry Operations Panel")

txBox = slide7.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(11.1), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "A high-end dark luxury telemetry interface that gives operations managers oversight over the ALE's operations."
p.font.name = "Segoe UI"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE
p.space_after = Pt(15)

dashboard_features = [
    ("Telemetry Data Visualization", "Recharts-powered interactive line and bar charts tracking revenue, dynamic weather metrics, and profit margins."),
    ("Live Transaction Feed", "Real-time updates showing block transactions, deployment hashes, inventory buying fees, and deposit logs."),
    ("AI Command Terminal", "Provides a terminal interface allowing developers to trigger mock outage states, change weather oracle data, or force CEO ticks manually."),
    ("Dynamic State Controls", "Switch state between online/offline operations instantly, demonstrating state synchronization between local SQLite and virtual blockchain nodes.")
]

for title, desc in dashboard_features:
    p_title = tf.add_paragraph()
    p_title.text = f"• {title}: "
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(13)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    p_title.space_before = Pt(6)

    p_desc = tf.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = "Segoe UI"
    p_desc.font.size = Pt(11)
    p_desc.font.color.rgb = TEXT_BODY
    p_desc.space_after = Pt(6)

add_footer(slide7, 7)

# ----------------- SLIDE 8: DYNAMIC PRICING ENGINE -----------------
slide8 = prs.slides.add_slide(slide_layout)
apply_background(slide8)
add_header(slide8, "Dynamic Pricing Engine: Maximizing Profit Margins")

# Left Box: Pricing Mechanics
txBox_left = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5))
tf_left = txBox_left.text_frame
tf_left.word_wrap = True

p = tf_left.paragraphs[0]
p.text = "How the Dynamic Engine Works"
p.font.name = "Segoe UI"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE
p.space_after = Pt(15)

pricing_bullets = [
    ("Weather Oracle Alignment", "If weather changes to 'Hot/Sunny', consumer demand index jumps up. The AI automatically increases the lemonade price per cup (e.g. ₹20-25)."),
    ("Rainy/Cloudy Mitigation", "In case of rainy weather (drop in traffic), the AI reduces prices to ₹8-10 to trigger volume sales and avoid ingredient expiration."),
    ("Inventory Overhead Calculations", "Calculates holding costs. If ice cubes are melting or inventory limits are near, prices drop to clear out materials."),
    ("Competitor & Cost Indexing", "Tracks wholesale cost changes. If wholesale lemon price spikes, the minimum price threshold shifts automatically.")
]

for title, desc in pricing_bullets:
    p_title = tf_left.add_paragraph()
    p_title.text = f"• {title}"
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(13)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    p_title.space_before = Pt(6)
    
    p_desc = tf_left.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = "Segoe UI"
    p_desc.font.size = Pt(10)
    p_desc.font.color.rgb = TEXT_BODY
    p_desc.space_after = Pt(4)

# Right Box: Large Statistics Callout
add_card(slide8, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5), "AI Engine Yield Optimization", border_color=ACCENT_GREEN)
txBox_right = slide8.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5.1), Inches(3.5))
tf_right = txBox_right.text_frame
tf_right.word_wrap = True

p_call = tf_right.paragraphs[0]
p_call.text = "During simulations, dynamic pricing adjustments outperformed static pricing schemes by increasing overall yield:"
p_call.font.name = "Segoe UI"
p_call.font.size = Pt(14)
p_call.font.color.rgb = TEXT_BODY
p_call.space_after = Pt(25)

p_stat = tf_right.add_paragraph()
p_stat.text = "+42%"
p_stat.font.name = "Segoe UI"
p_stat.font.size = Pt(54)
p_stat.font.bold = True
p_stat.font.color.rgb = ACCENT_GREEN
p_stat.alignment = PP_ALIGN.CENTER

p_stat_lbl = tf_right.add_paragraph()
p_stat_lbl.text = "Increase in overall business profitability margins"
p_stat_lbl.font.name = "Segoe UI"
p_stat_lbl.font.size = Pt(11)
p_stat_lbl.font.color.rgb = TEXT_MUTED
p_stat_lbl.alignment = PP_ALIGN.CENTER

add_footer(slide8, 8)

# ----------------- SLIDE 9: DECENTRALIZED WORKFORCE -----------------
slide9 = prs.slides.add_slide(slide_layout)
apply_background(slide9)
add_header(slide9, "Decentralized Workforce: Autonomous Recruitment")

# 3 Columns for Worker Journey
worker_steps = [
    ("1. Job Postings", "When operations begin or scaling is required, the AI posts jobs to the on-chain board with payment rates (e.g. 'Produce Lemonade', Rate: ₹200).", Inches(0.8)),
    ("2. Identity Applications", "Gig workers apply by linking their web3 wallet and providing identity info. AI registers them in the contract database.", Inches(4.8)),
    ("3. Automated Payroll", "Once the task is submitted, the AI verifies the job status and issues a direct payout transaction. Zero billing paperwork required.", Inches(8.8))
]

for title, desc, left in worker_steps:
    add_card(slide9, left, Inches(1.8), Inches(3.7), Inches(4.5), title, border_color=ACCENT_BLUE)
    txBox = slide9.shapes.add_textbox(left + Inches(0.25), Inches(2.6), Inches(3.2), Inches(3.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.name = "Segoe UI"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_BODY
    p.space_after = Pt(20)

    p_status = tf.add_paragraph()
    p_status.text = "VERIFIED CONTRACT"
    p_status.font.name = "Segoe UI"
    p_status.font.size = Pt(9)
    p_status.font.bold = True
    p_status.font.color.rgb = ACCENT_BLUE

add_footer(slide9, 9)

# ----------------- SLIDE 10: THE BLOCKCHAIN TREASURY -----------------
slide10 = prs.slides.add_slide(slide_layout)
apply_background(slide10)
add_header(slide10, "The Blockchain Treasury: Auditable Ledger")

# Left Column: Treasury Rules
txBox_left = slide10.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5))
tf_left = txBox_left.text_frame
tf_left.word_wrap = True

p = tf_left.paragraphs[0]
p.text = "Decentralized Capital Management"
p.font.name = "Segoe UI"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE
p.space_after = Pt(15)

treasury_points = [
    ("Non-Custodial Account Structures", "Capital is held in transparent blockchain wallets, meaning shareholders can track every transaction in real-time."),
    ("Multi-Sig Operational Safe", "Large expenses (like inventory restocks above ₹1000) require secondary oracle validations to prevent AI exploitation."),
    ("Automated Gas & Fee Handling", "Transactions run with optimized gas pricing estimations, keeping chain interaction costs at absolute minimums."),
    ("One-Click Deposit & Withdrawals", "Managers can deposit liquidity to seed the business or withdraw profits through verified administrative keys.")
]

for title, desc in treasury_points:
    p_title = tf_left.add_paragraph()
    p_title.text = f"✔ {title}"
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(13)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    p_title.space_before = Pt(6)
    
    p_desc = tf_left.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = "Segoe UI"
    p_desc.font.size = Pt(10)
    p_desc.font.color.rgb = TEXT_BODY
    p_desc.space_after = Pt(4)

# Right Column: Security ledger info
add_card(slide10, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5), "On-Chain Audit Log", border_color=ACCENT_PURPLE)
txBox_right = slide10.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5.1), Inches(3.5))
tf_right = txBox_right.text_frame
tf_right.word_wrap = True

p_aud = tf_right.paragraphs[0]
p_aud.text = "Audit Trails & Security Specs:"
p_aud.font.name = "Segoe UI"
p_aud.font.size = Pt(14)
p_aud.font.bold = True
p_aud.font.color.rgb = TEXT_WHITE
p_aud.space_after = Pt(15)

audit_bullets = [
    ("Transaction Hashing", "Every purchase, sale, and wage payout generates a unique transaction hash (e.g. 0x8eCe6838D...) for public indexing."),
    ("Immutable Ledger Records", "Database transactions sync to blockchain blocks, making it impossible to backdate receipts or alter profit records."),
    ("Real-Time Balance Verification", "Asset balance checks run continuously, cross-referencing local SQLite states with contract wallet values.")
]

for title, desc in audit_bullets:
    p_a_title = tf_right.add_paragraph()
    p_a_title.text = f"• {title}"
    p_a_title.font.name = "Segoe UI"
    p_a_title.font.size = Pt(11)
    p_a_title.font.bold = True
    p_a_title.font.color.rgb = ACCENT_PURPLE
    p_a_title.space_before = Pt(6)
    
    p_a_desc = tf_right.add_paragraph()
    p_a_desc.text = desc
    p_a_desc.font.name = "Segoe UI"
    p_a_desc.font.size = Pt(9.5)
    p_a_desc.font.color.rgb = TEXT_BODY
    p_a_desc.space_after = Pt(4)

add_footer(slide10, 10)

# ----------------- SLIDE 11: MARKET ORACLE SYSTEM -----------------
slide11 = prs.slides.add_slide(slide_layout)
apply_background(slide11)
add_header(slide11, "Market Oracle System: Bridge to the Real World")

add_card(slide11, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5), "Connecting Real-World Telemetry to Smart Contracts")

txBox = slide11.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(11.1), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Autonomous legal entities cannot operate in silos. ATLAS utilizes a data bridge ('Oracle') to bring real-world metrics onto the ledger."
p.font.name = "Segoe UI"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE
p.space_after = Pt(15)

oracle_bullets = [
    ("Weather Data Streams", "Monitors temperature, precipitation, and humidity index to compute local customer traffic coefficients."),
    ("Wholesale Pricing Feeds", "Ingests live material costs (e.g. lemons, sugar) to set maximum budget thresholds for inventory ordering."),
    ("Demand Volatility Analysis", "Evaluates dynamic consumer interest to balance price optimization against supply exhaustion risks."),
    ("Verification Mechanics", "Multi-source oracle checks ensure that individual bad sensors or feed disruptions cannot manipulate the AI's core pricing logic.")
]

for title, desc in oracle_bullets:
    p_title = tf.add_paragraph()
    p_title.text = f"✔ {title}: "
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(13)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    p_title.space_before = Pt(6)

    p_desc = tf.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = "Segoe UI"
    p_desc.font.size = Pt(11)
    p_desc.font.color.rgb = TEXT_BODY
    p_desc.space_after = Pt(6)

add_footer(slide11, 11)

# ----------------- SLIDE 12: SECURITY & GUARDRAILS -----------------
slide12 = prs.slides.add_slide(slide_layout)
apply_background(slide12)
add_header(slide12, "Security & Guardrails: Preventing AI Exploitability")

# Left Column: Risk vectors
txBox_left = slide12.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5))
tf_left = txBox_left.text_frame
tf_left.word_wrap = True

p = tf_left.paragraphs[0]
p.text = "Mitigating Autonomous Risks"
p.font.name = "Segoe UI"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE
p.space_after = Pt(15)

risk_points = [
    ("Limit & Budget Caps", "Hard limits on transactions (e.g. max ₹1500 per day for restocks) prevent prompt-injection attacks from draining treasury funds."),
    ("Sandboxed Decision Testing", "Before pushing changes to the live SQLite db or blockchain ledger, decisions run through a Dry-Run Validator to check for anomalies."),
    ("Fallback Control Protocol", "If any exception occurs in the reasoning pipeline, the system automatically triggers a Safe Fallback mode, reverting variables to default states.")
]

for title, desc in risk_points:
    p_title = tf_left.add_paragraph()
    p_title.text = f"• {title}"
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(13)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    p_title.space_before = Pt(6)
    
    p_desc = tf_left.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = "Segoe UI"
    p_desc.font.size = Pt(10.5)
    p_desc.font.color.rgb = TEXT_BODY
    p_desc.space_after = Pt(4)

# Right Column: Big Alert Card
add_card(slide12, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5), "Fail-Safe Governance", border_color=ACCENT_GREEN)
txBox_right = slide12.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5.1), Inches(3.5))
tf_right = txBox_right.text_frame
tf_right.word_wrap = True

p_alert = tf_right.paragraphs[0]
p_alert.text = "Decentralized businesses require bulletproof security against AI drift or prompt tampering."
p_alert.font.name = "Segoe UI"
p_alert.font.size = Pt(14)
p_alert.font.color.rgb = TEXT_BODY
p_alert.space_after = Pt(20)

alerts = [
    ("Administrative Sign-offs", "Certain high-value decisions (e.g., wallet withdrawals) generate a verification link requiring admin approval, maintaining human-in-the-loop oversight."),
    ("Immutable Logging", "Every prompt sent to Claude, and every raw response received, is stored in a read-only ActivityLog, creating an unalterable history ledger.")
]

for title, desc in alerts:
    p_a = tf_right.add_paragraph()
    p_a.text = f"✔ {title}"
    p_a.font.name = "Segoe UI"
    p_a.font.size = Pt(12)
    p_a.font.bold = True
    p_a.font.color.rgb = ACCENT_GREEN
    p_a.space_before = Pt(8)
    
    p_ad = tf_right.add_paragraph()
    p_ad.text = desc
    p_ad.font.name = "Segoe UI"
    p_ad.font.size = Pt(10)
    p_ad.font.color.rgb = TEXT_BODY
    p_ad.space_after = Pt(4)

add_footer(slide12, 12)

# ----------------- SLIDE 13: COMPETITIVE ADVANTAGE -----------------
slide13 = prs.slides.add_slide(slide_layout)
apply_background(slide13)
add_header(slide13, "ATLAS AI: Competitive Edge & Market Value")

# 3 columns comparative cards
comp_data = [
    ("Traditional Enterprise", "• High corporate management overhead\n• Manual billing and payroll processes\n• Slow reaction to localized market changes\n• Closed internal records with opaque audits\n• Relies on slow, manual vendor contract enforcement", Inches(0.8), TEXT_MUTED),
    ("AI Co-Pilot Solutions", "• AI assists human managers\n• Manual actions still required to execute transactions\n• High cost of human error remains\n• Limited real-time autonomy\n• Lacks direct ledger/wallet integration", Inches(4.8), ACCENT_PURPLE),
    ("ATLAS ALE (Our Solution)", "• Zero management overhead costs\n• Real-time, 24/7 automated operations\n• Programmatic pricing based on live Oracles\n• Pure blockchain transparency & auditing\n• Automated smart contract workforce payouts", Inches(8.8), ACCENT_GREEN)
]

for title, list_text, left, color in comp_data:
    add_card(slide13, left, Inches(1.8), Inches(3.7), Inches(4.5), title, border_color=color)
    txBox = slide13.shapes.add_textbox(left + Inches(0.2), Inches(2.5), Inches(3.3), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for line in list_text.split('\n'):
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        p.text = line
        p.font.name = "Segoe UI"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_BODY
        p.space_after = Pt(6)

add_footer(slide13, 13)

# ----------------- SLIDE 14: FUTURE ROADMAP -----------------
slide14 = prs.slides.add_slide(slide_layout)
apply_background(slide14)
add_header(slide14, "Future Milestones & Scaling")

# Roadmap Grid (4 Cards)
roadmap_milestones = [
    ("Q3 2026: Multi-Chain Deployments", "Support EVM chains, Solana, and Layer 2s (Arbitrum/Optimism) for gas optimization and microtransaction scaling.", Inches(0.8)),
    ("Q4 2026: Sandboxed Previews", "Provide isolated container sandboxes to dry-run AI decisions, confirming there are no smart contract regressions before pushing live.", Inches(3.8)),
    ("Q1 2027: ChatOps Integrations", "Provide deep Slack, Discord, and Telegram integrations to notify operators about critical ledger changes or worker applications.", Inches(6.8)),
    ("Q2 2027: LLM Orchestrator Framework", "Allow developers to deploy modular custom agents (e.g. Finance Agent, HR Agent) with customizable risk thresholds.", Inches(9.8))
]

for title, desc, left in roadmap_milestones:
    add_card(slide14, left, Inches(1.8), Inches(2.7), Inches(4.5), title, border_color=ACCENT_BLUE)
    txBox = slide14.shapes.add_textbox(left + Inches(0.2), Inches(2.5), Inches(2.3), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.name = "Segoe UI"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_BODY
    p.space_after = Pt(15)

    p_tag = tf.add_paragraph()
    p_tag.text = "PLANNED PROTOCOL"
    p_tag.font.name = "Segoe UI"
    p_tag.font.size = Pt(9)
    p_tag.font.bold = True
    p_tag.font.color.rgb = ACCENT_BLUE

add_footer(slide14, 14)

# ----------------- SLIDE 15: CONCLUSION & Q&A -----------------
slide15 = prs.slides.add_slide(slide_layout)
apply_background(slide15)
add_header(slide15, "Join the Decentralized Future", "CONTACT // Q&A")

add_card(slide15, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5), "ATLAS AI: Build Autonomous. Operate Everywhere.", border_color=ACCENT_PURPLE)

txBox = slide15.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(11.1), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "Thank You!"
p1.font.name = "Segoe UI"
p1.font.size = Pt(36)
p1.font.bold = True
p1.font.color.rgb = TEXT_WHITE
p1.space_after = Pt(10)

p2 = tf.add_paragraph()
p2.text = "ATLAS AI represents the paradigm shift in corporate operations: moving from human bureaucratic friction to automated, trustless, and intelligent Web3 command systems."
p2.font.name = "Segoe UI"
p2.font.size = Pt(14)
p2.font.color.rgb = TEXT_BODY
p2.space_after = Pt(25)

p3 = tf.add_paragraph()
p3.text = "GitHub Repository: github.com/aayush-mistry/ATLAS-AI"
p3.font.name = "Segoe UI"
p3.font.size = Pt(13)
p3.font.bold = True
p3.font.color.rgb = ACCENT_BLUE
p3.space_after = Pt(8)

p4 = tf.add_paragraph()
p4.text = "Contact Team: meridian-devs@atlas.ai"
p4.font.name = "Segoe UI"
p4.font.size = Pt(13)
p4.font.bold = True
p4.font.color.rgb = ACCENT_PURPLE

add_footer(slide15, 15)

# Save presentation
prs.save("ATLAS_AI_Pitch_Deck.pptx")
print("Presentation created successfully as 'ATLAS_AI_Pitch_Deck.pptx'!")
